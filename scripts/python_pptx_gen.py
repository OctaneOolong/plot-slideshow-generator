from pptx import Presentation
from pptx.util import Inches
import os, sys, subprocess
import time
from pathlib import Path

def slide_image_adder(img_path_, prs_):
# get the image

    img_path = img_path_

    # create a blank slide that the image will be placed on
    blank_slide_layout = prs_.slide_layouts[6]
    slide = prs_.slides.add_slide(blank_slide_layout)

    # define the image offset
    left = top = Inches(0)

    # add the image
    pic = slide.shapes.add_picture(img_path, left, top)

    return prs_

def python_pptx_gen(images_dir_list, delete_input_files = bool):
    """
    pass the image files as a path iterable, i.e. Path.glob().

    takes a list of images, adds one image per slide, for as many images as are provided, returns the presentation object. 

    Currently the images need to be saved to file, I guess as .png atm. 

    Option to clean up afterwards:
    delete_input_files: 'y' or 'n'. 'y' will delete all the input files after the presentation is produced, cleaning up the directory. 'n' will continue, leaving the files intact. 
    """
    
    # Produce the presentation default format
    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)

    dir_path = Path.cwd()

    for image_path in images_dir_list:
        prs = slide_image_adder(Path(image_path).name, prs)
        if delete_input_files == 'y':
            image_path.unlink()

    # define and save the output filename

    return prs

if __name__ == "__main__":
    python_pptx_gen()